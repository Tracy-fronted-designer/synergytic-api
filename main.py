from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import Optional
import httpx
import io

from pptx import Presentation
from pptx.util import Pt
from openpyxl import Workbook

app = FastAPI(
    title="Synergytic API",
    description="信瑞企管 AI 企劃書自動化系統 — PPT 產出與 Excel 索引建立",
    version="1.0.0",
)

GRAPH_BASE = "https://graph.microsoft.com/v1.0"
ONEDRIVE_ROOT_ID = "42B13569E5D6C5E7!533"


# ── Models ──────────────────────────────────────────────────────────────────

class ScanProposalsRequest(BaseModel):
    access_token: str
    folder_id: Optional[str] = ONEDRIVE_ROOT_ID

class GenerateIndexRequest(BaseModel):
    access_token: str
    proposals: list[dict]  # 由 /scan-proposals 回傳的清單

class Section(BaseModel):
    title: str
    content: str

class GeneratePptxRequest(BaseModel):
    client_name: str
    course_title: str
    instructor: str
    sections: list[Section]


# ── Endpoints ────────────────────────────────────────────────────────────────

@app.get("/health", summary="健康檢查")
def health():
    return {"status": "ok"}


@app.post("/scan-proposals", summary="掃描 OneDrive 所有企劃書")
async def scan_proposals(req: ScanProposalsRequest):
    """
    列出指定 OneDrive 資料夾下所有子資料夾，
    再進入每個子資料夾找出 (信瑞企劃書) 開頭的 PPT/PDF 檔案。
    回傳結構化清單供後續建立索引用。
    """
    headers = {"Authorization": f"Bearer {req.access_token}"}
    results = []

    async with httpx.AsyncClient() as client:
        # 取得所有客戶資料夾
        folders_url = f"{GRAPH_BASE}/me/drive/items/{req.folder_id}/children"
        resp = await client.get(folders_url, headers=headers)
        if resp.status_code != 200:
            raise HTTPException(status_code=resp.status_code, detail=resp.text)

        folders = resp.json().get("value", [])

        for folder in folders:
            if folder.get("folder") is None:
                continue

            folder_id = folder["id"]
            folder_name = folder["name"]

            # 取得資料夾內檔案
            files_url = f"{GRAPH_BASE}/me/drive/items/{folder_id}/children"
            files_resp = await client.get(files_url, headers=headers)
            if files_resp.status_code != 200:
                continue

            for f in files_resp.json().get("value", []):
                name: str = f.get("name", "")
                if not name.startswith("(信瑞企劃書)"):
                    continue
                ext = name.rsplit(".", 1)[-1].lower()
                if ext not in ("ppt", "pptx", "pdf"):
                    continue

                results.append({
                    "client_name": folder_name,
                    "file_name": name,
                    "file_id": f["id"],
                    "created": f.get("createdDateTime", ""),
                    "web_url": f.get("webUrl", ""),
                })

    return {"count": len(results), "proposals": results}


@app.post("/generate-index", summary="產出 Excel 索引檔（Base64）")
async def generate_index(req: GenerateIndexRequest):
    """
    接收企劃書清單（含 AI 解析好的欄位），
    產出 Excel 索引檔並以 Base64 回傳，供 n8n 存到 OneDrive。
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "企劃書索引"

    headers = [
        "客戶名稱", "產業別", "課程主題", "痛點標籤", "適合對象",
        "講師", "課程天數", "課程模式", "摘要", "檔案名稱", "檔案ID", "建立日期",
    ]
    ws.append(headers)

    for p in req.proposals:
        ws.append([
            p.get("client_name", ""),
            p.get("industry", ""),
            p.get("course_topic", ""),
            p.get("pain_points", ""),
            p.get("target_audience", ""),
            p.get("instructor", ""),
            p.get("days", ""),
            p.get("mode", ""),
            p.get("summary", ""),
            p.get("file_name", ""),
            p.get("file_id", ""),
            p.get("created", ""),
        ])

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    import base64
    b64 = base64.b64encode(buf.read()).decode()
    return {"filename": "企劃書索引.xlsx", "content_base64": b64}


@app.post("/generate-pptx", summary="產出企劃書 PPT（Base64）")
async def generate_pptx(req: GeneratePptxRequest):
    """
    依傳入的客戶名稱、課程標題、講師、內容段落，
    產出 PPT 並以 Base64 回傳，供 n8n 存到 OneDrive。
    """
    prs = Presentation()
    slide_layout_title = prs.slide_layouts[0]   # Title Slide
    slide_layout_content = prs.slide_layouts[1]  # Title and Content

    # 封面
    slide = prs.slides.add_slide(slide_layout_title)
    slide.shapes.title.text = req.course_title
    slide.placeholders[1].text = f"客戶：{req.client_name}　｜　講師：{req.instructor}"

    # 內容頁
    for section in req.sections:
        slide = prs.slides.add_slide(slide_layout_content)
        slide.shapes.title.text = section.title
        tf = slide.placeholders[1].text_frame
        tf.text = section.content
        tf.paragraphs[0].font.size = Pt(16)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    import base64
    b64 = base64.b64encode(buf.read()).decode()
    filename = f"{req.client_name}_{req.course_title}.pptx"
    return {"filename": filename, "content_base64": b64}
